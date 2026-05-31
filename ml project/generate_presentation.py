from pptx import Presentation
from pptx.util import Inches, Pt

slides = [
    {
        'title': 'Retail Demand Forecasting Project',
        'content': 'A complete end-to-end solution for retail sales forecasting, AI insights, and deployment.'
    },
    {
        'title': 'Project Objective',
        'content': 'Forecast daily retail sales and provide actionable business recommendations using machine learning and GenAI.'
    },
    {
        'title': 'Data Source',
        'content': 'Rossmann Store Sales dataset with historical sales, promotions, holidays, store metadata, and competition information.'
    },
    {
        'title': 'Key Features',
        'content': 'Store ID, DayOfWeek, Promo, StateHoliday, StoreType, Assortment, CompetitionDistance, Year, Month, Day, WeekOfYear.'
    },
    {
        'title': 'Data Preprocessing',
        'content': 'Cleaned missing values, encoded categorical features, extracted date features, and filtered open store sales.'
    },
    {
        'title': 'Exploratory Data Analysis',
        'content': 'Generated sales distribution, promo impact, seasonality, and correlation visualizations.'
    },
    {
        'title': 'Modeling Approach',
        'content': 'Trained Linear Regression and Random Forest models to compare interpretability and prediction performance.'
    },
    {
        'title': 'Evaluation Metrics',
        'content': 'Used RMSE and MAE to measure forecast accuracy and compare model performance.'
    },
    {
        'title': 'Streamlit Dashboard',
        'content': 'Interactive dashboard for what-if scenarios, model predictions, and AI narrative explanations.'
    },
    {
        'title': 'GenAI Integration',
        'content': 'Anthropic Claude API provides strategic insights, chat-based analysis, and anomaly explanations.'
    },
    {
        'title': 'Deployment Assets',
        'content': 'Includes Dockerfile, Procfile, requirements, and deployment guide for local and cloud hosting.'
    },
    {
        'title': 'Deployment Options',
        'content': 'Local Python, Docker container, Streamlit Cloud, or cloud PaaS using the provided configuration.'
    },
    {
        'title': 'Code Structure',
        'content': 'Modular design with data_loader, preprocessing, eda, model training, anomaly detection, and AI engine modules.'
    },
    {
        'title': 'Business Impact',
        'content': 'Helps retail teams improve inventory planning, staffing, and promotional decisions through demand forecasts.'
    },
    {
        'title': 'Conclusion',
        'content': 'The project delivers an end-to-end retail demand forecasting system with analytics, insights, and deployment readiness.'
    },
]

presentation = Presentation()
for i, slide_data in enumerate(slides):
    slide_layout = presentation.slide_layouts[0] if i == 0 else presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = slide_data['title']

    if len(slide.shapes.placeholders) > 1:
        body = slide.shapes.placeholders[1]
        text_frame = body.text_frame
        text_frame.text = slide_data['content']
    else:
        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)
        height = Inches(4.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = slide_data['content']

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)

presentation.save('Retail_Demand_Forecasting_Presentation.pptx')
print('Presentation generated: Retail_Demand_Forecasting_Presentation.pptx')
